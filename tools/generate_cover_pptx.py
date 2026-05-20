"""
Capa do ebook: High Protein Cookbook for Beginners
Abordagem 2: python-pptx
Manifesto visual: Honest Kitchen
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

OUT = r"C:\Users\gardi\OneDrive\Documentos\BRAIN\.tmp\cover_pptx.pptx"

# Paleta
BG       = RGBColor(0xFA, 0xF8, 0xF4)
CHARCOAL = RGBColor(0x2C, 0x2C, 0x2C)
TERRA    = RGBColor(0xC4, 0x62, 0x2D)
SAGE     = RGBColor(0x6B, 0x8F, 0x71)
DIVIDER  = RGBColor(0xE8, 0xE0, 0xD5)
LIGHT_SAGE = RGBColor(0xD4, 0xE0, 0xD8)
GRAY     = RGBColor(0x9B, 0x9B, 0x9B)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)

# KDP cover ratio ~1:1.6 → 6x9.6 inches
prs = Presentation()
prs.slide_width  = Inches(6)
prs.slide_height = Inches(9.6)

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

W = prs.slide_width
H = prs.slide_height

def add_rect(slide, x, y, w, h, fill_color, border_color=None, border_width=0):
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_width)
    else:
        shape.line.fill.background()
    return shape

def add_textbox(slide, text, x, y, w, h, font_name, font_size, font_color,
                bold=False, align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.color.rgb = font_color
    run.font.bold = bold
    run.font.italic = italic
    return txBox

# --- BACKGROUND ---
add_rect(slide, 0, 0, W, H, BG)

# --- FAIXA TERRACOTA ESQUERDA ---
add_rect(slide, 0, 0, Inches(0.08), H, TERRA)

# --- FOTO PLACEHOLDER (lado esquerdo) ---
photo_w = int(W * 0.52)
photo_h = int(H * 0.62)
photo_y = int(H * 0.19)
add_rect(slide, 0, photo_y, photo_w, photo_h, LIGHT_SAGE)

# Texto no placeholder
add_textbox(slide, "[ food photo ]",
            Inches(0.1), photo_y + photo_h//2 - Inches(0.2),
            photo_w - Inches(0.2), Inches(0.4),
            "Arial", 8, SAGE, align=PP_ALIGN.CENTER)

# --- LINHA DIVISÓRIA VERTICAL ---
line = slide.shapes.add_shape(1,
    photo_w + Inches(0.1),
    int(H * 0.08),
    Pt(1),
    int(H * 0.84))
line.fill.solid()
line.fill.fore_color.rgb = DIVIDER
line.line.fill.background()

# --- ÁREA DE TEXTO DIREITA ---
tx = photo_w + Inches(0.22)
tw = W - tx - Inches(0.25)

# EST. 2026
add_textbox(slide, "EST. 2026",
            tx, Inches(0.4), tw, Inches(0.25),
            "Arial", 7, TERRA, bold=False)

# COOKBOOK tag
add_textbox(slide, "COOKBOOK",
            tx, Inches(0.68), tw, Inches(0.25),
            "Arial", 7, SAGE, bold=True)

# Linha sage sob tag
line2 = slide.shapes.add_shape(1, tx, Inches(0.95), int(tw * 0.9), Pt(1))
line2.fill.solid()
line2.fill.fore_color.rgb = SAGE
line2.line.fill.background()

# HIGH
add_textbox(slide, "HIGH",
            tx, Inches(1.05), tw, Inches(0.55),
            "Georgia", 30, CHARCOAL, bold=True)

# PROTEIN (terracota)
add_textbox(slide, "PROTEIN",
            tx, Inches(1.55), tw, Inches(0.55),
            "Georgia", 30, TERRA, bold=True)

# COOKBOOK
add_textbox(slide, "COOKBOOK",
            tx, Inches(2.05), tw, Inches(0.45),
            "Georgia", 22, CHARCOAL, bold=True)

# Linha terracota
line3 = slide.shapes.add_shape(1, tx, Inches(2.55), int(tw * 0.85), Pt(2))
line3.fill.solid()
line3.fill.fore_color.rgb = TERRA
line3.line.fill.background()

# FOR BEGINNERS
add_textbox(slide, "FOR BEGINNERS",
            tx, Inches(2.65), tw, Inches(0.3),
            "Arial", 9, CHARCOAL, bold=False)

# Subtítulo
add_textbox(slide, "50 Easy Recipes to Build\nMuscle & Eat Well",
            tx, Inches(3.05), tw, Inches(0.6),
            "Arial", 8, RGBColor(0x6B,0x6B,0x6B))

# Círculos de benefícios
benefits = [("HIGH\nPROTEIN", 0), ("30 MIN\nMEALS", 1), ("BEGINNER\nFRIENDLY", 2)]
circle_size = Inches(0.55)
circle_y = int(H * 0.60)
spacing = int((tw - circle_size) / 2)

for label, i in benefits:
    cx = tx + i * (circle_size + spacing // 2)
    shape = slide.shapes.add_shape(9, cx, circle_y, circle_size, circle_size)  # oval
    shape.fill.solid()
    shape.fill.fore_color.rgb = SAGE
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.name = "Arial"
    run.font.size = Pt(5.5)
    run.font.bold = True
    run.font.color.rgb = WHITE

# Linha rodapé
line4 = slide.shapes.add_shape(1, tx, int(H * 0.88), int(tw * 0.9), Pt(1))
line4.fill.solid()
line4.fill.fore_color.rgb = SAGE
line4.line.fill.background()

# Rodapé texto
add_textbox(slide, "High Protein Cookbook for Beginners",
            tx, int(H * 0.89), tw, Inches(0.25),
            "Arial", 6.5, GRAY)

prs.save(OUT)
print(f"PPTX gerado: {OUT}")
